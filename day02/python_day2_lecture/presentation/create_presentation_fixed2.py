"""
Create PowerPoint presentation for Python Day 2 lecture covering Functions, Loops, and Data Structures
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    """Add a title slide with the given title and subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_section_slide(prs, title):
    """Add a section slide with the given title."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_content_slide(prs, title, content, code=None):
    """Add a content slide with the given title, content, and optional code."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        
        # Check if this is a bullet point
        if line.startswith('•'):
            p.level = 1
    
    # Add code if provided
    if code:
        left = Inches(0.5)
        top = Inches(4)
        width = Inches(9)
        height = Inches(2.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.name = 'Courier New'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        # Add a light gray background to the code box
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide

def create_presentation():
    """Create the PowerPoint presentation for Python Day 2 lecture."""
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, 
                   "Python Functions, Loops, and Data Structures", 
                   "Day 2: Building Blocks of Python Programming")
    
    # Introduction slide
    add_content_slide(prs, 
                     "Today's Topics", 
                     [
                         "• Functions: definition, parameters, return types",
                         "• Loops: While and For loops",
                         "• Data Structures: Lists, Tuples, Sets, and Dictionaries",
                         "",
                         "By the end of this session, you'll be able to:",
                         "• Create and use functions to organize code",
                         "• Use loops to repeat code efficiently",
                         "• Work with different data structures to store and manage information",
                         "• Apply these concepts in a real-life Student Grading System"
                     ])
    
    # Functions section
    add_section_slide(prs, "Functions in Python")
    
    add_content_slide(prs, 
                     "What are Functions?", 
                     [
                         "• Reusable blocks of code designed to perform a specific task",
                         "• Help organize code and make it reusable",
                         "• Break down complex problems into smaller, manageable parts",
                         "",
                         "Benefits of using functions:",
                         "• Code reusability",
                         "• Better organization",
                         "• Easier debugging",
                         "• Improved readability"
                     ])
    
    add_content_slide(prs, 
                     "Defining Functions", 
                     [
                         "• Use the def keyword followed by function name and parentheses",
                         "• Function body is indented",
                         "• Call the function by using its name followed by parentheses"
                     ],
                     code="# Defining a function\ndef greet():\n    print(\"Hello, World!\")\n\n# Calling the function\ngreet()  # Output: Hello, World!\n\n# Function with parameters\ndef greet_person(name):\n    print(f\"Hello, {name}!\")\n\ngreet_person(\"Alice\")  # Output: Hello, Alice!")
    
    add_content_slide(prs, 
                     "Function Parameters", 
                     [
                         "• Parameters allow functions to receive input values",
                         "• Multiple parameters are separated by commas",
                         "• Default parameters provide fallback values",
                         "• Keyword arguments let you specify which parameter gets which value"
                     ],
                     code="# Multiple parameters\ndef describe_person(name, age):\n    print(f\"{name} is {age} years old.\")\n\ndescribe_person(\"Charlie\", 25)  # Output: Charlie is 25 years old.\n\n# Default parameter values\ndef greet_with_message(name, message=\"Good morning\"):\n    print(f\"{message}, {name}!\")\n\ngreet_with_message(\"David\")  # Output: Good morning, David!\ngreet_with_message(\"Eve\", \"Good evening\")  # Output: Good evening, Eve!")
    
    add_content_slide(prs, 
                     "Return Values", 
                     [
                         "• Functions can return values using the return statement",
                         "• Return values can be assigned to variables",
                         "• Functions can return multiple values as a tuple",
                         "• A function without a return statement returns None"
                     ],
                     code="# Function that returns a value\ndef add_numbers(a, b):\n    return a + b\n\nsum_result = add_numbers(5, 3)\nprint(sum_result)  # Output: 8\n\n# Function that returns multiple values\ndef get_dimensions():\n    return 500, 300\n\nwidth, height = get_dimensions()\nprint(f\"Width: {width}, Height: {height}\")  # Output: Width: 500, Height: 300")
    
    add_content_slide(prs, 
                     "Variable Scope", 
                     [
                         "• Variables defined inside a function have local scope",
                         "• Variables defined outside functions have global scope",
                         "• Use the global keyword to modify global variables inside a function"
                     ],
                     code="# Local scope\ndef calculate():\n    local_var = 10\n    print(local_var)  # This works\n\ncalculate()\n# print(local_var)  # This would cause an error\n\n# Global scope\nglobal_var = 20\n\ndef print_global():\n    print(global_var)  # This works\n\nprint_global()  # Output: 20\n\n# Modifying global variables\ncounter = 0\n\ndef increment():\n    global counter\n    counter += 1\n    print(counter)\n\nincrement()  # Output: 1")
    
    add_content_slide(prs, 
                     "Lambda Functions", 
                     [
                         "• Small, anonymous functions defined with the lambda keyword",
                         "• Can take any number of arguments but only have one expression",
                         "• Useful for short, simple functions that are used only once",
                         "• Often used with functions like map(), filter(), and sorted()"
                     ],
                     code="# Regular function vs. lambda function\ndef double(x):\n    return x * 2\n\ndouble_lambda = lambda x: x * 2\n\nprint(double(5))        # Output: 10\nprint(double_lambda(5)) # Output: 10\n\n# Using lambda with filter()\nnumbers = [1, 2, 3, 4, 5]\neven_numbers = list(filter(lambda x: x % 2 == 0, numbers))\nprint(even_numbers)  # Output: [2, 4]")
    
    # Loops section
    add_section_slide(prs, "Loops in Python")
    
    add_content_slide(prs, 
                     "What are Loops?", 
                     [
                         "• Loops allow you to execute a block of code multiple times",
                         "• Avoid repetitive code",
                         "• Two main types of loops in Python:",
                         "  - while loops: Execute as long as a condition is true",
                         "  - for loops: Iterate over a sequence",
                         "",
                         "• Special statements:",
                         "  - break: Exit the loop",
                         "  - continue: Skip to the next iteration"
                     ])
    
    add_content_slide(prs, 
                     "While Loops", 
                     [
                         "• Execute a block of code as long as a condition is true",
                         "• Condition is checked before each iteration",
                         "• Must include a way to exit the loop (avoid infinite loops)",
                         "• Can use break to exit prematurely",
                         "• Can use continue to skip the current iteration"
                     ],
                     code="# Basic while loop\ncount = 1\nwhile count <= 5:\n    print(count)\n    count += 1\n# Output: 1 2 3 4 5\n\n# Using break\ncount = 1\nwhile True:\n    print(count)\n    count += 1\n    if count > 5:\n        break\n# Output: 1 2 3 4 5\n\n# Using continue\ncount = 0\nwhile count < 10:\n    count += 1\n    if count % 2 == 0:  # If count is even\n        continue        # Skip the rest of this iteration\n    print(count)\n# Output: 1 3 5 7 9")
    
    add_content_slide(prs, 
                     "For Loops", 
                     [
                         "• Iterate over a sequence (list, tuple, string, etc.)",
                         "• More concise than while loops for many tasks",
                         "• The range() function generates a sequence of numbers",
                         "• Can nest loops inside other loops",
                         "• Can have an else clause that executes when the loop completes normally"
                     ],
                     code="# Looping through a list\nfruits = [\"apple\", \"banana\", \"cherry\"]\nfor fruit in fruits:\n    print(fruit)\n# Output: apple banana cherry\n\n# Using range()\nfor i in range(5):  # 0 to 4\n    print(i)\n# Output: 0 1 2 3 4\n\n# Nested loops\nfor i in range(1, 3):\n    for j in range(1, 3):\n        print(f\"{i} x {j} = {i * j}\")\n# Output:\n# 1 x 1 = 1\n# 1 x 2 = 2\n# 2 x 1 = 2\n# 2 x 2 = 4")
    
    add_content_slide(prs, 
                     "Loop Examples", 
                     [
                         "• Calculating factorial",
                         "• Finding prime numbers",
                         "• Processing items in a collection",
                         "• Building patterns"
                     ],
                     code="# Calculate factorial\ndef factorial(n):\n    result = 1\n    for i in range(1, n + 1):\n        result *= i\n    return result\n\nprint(factorial(5))  # Output: 120\n\n# Print a simple pattern\nrows = 5\nfor i in range(1, rows + 1):\n    print('*' * i)\n# Output:\n# *\n# **\n# ***\n# ****\n# *****")
    
    # Data Structures section
    add_section_slide(prs, "Data Structures in Python")
    
    add_content_slide(prs, 
                     "Python Data Structures", 
                     [
                         "• Data structures help organize and store data efficiently",
                         "• Python has several built-in data structures:",
                         "  - Lists: Ordered, mutable collections",
                         "  - Tuples: Ordered, immutable collections",
                         "  - Sets: Unordered collections of unique elements",
                         "  - Dictionaries: Key-value pairs",
                         "",
                         "• Each has specific use cases and advantages"
                     ])
    
    # Lists
    add_content_slide(prs, 
                     "Lists", 
                     [
                         "• Ordered, mutable collections",
                         "• Can contain elements of different types",
                         "• Created with square brackets []",
                         "• Zero-indexed (first element is at index 0)",
                         "• Support slicing, concatenation, and repetition"
                     ],
                     code="# Creating lists\nempty_list = []\nnumbers = [1, 2, 3, 4, 5]\nmixed = [1, \"hello\", 3.14, True]\n\n# Accessing elements\nfruits = [\"apple\", \"banana\", \"cherry\", \"date\"]\nprint(fruits[0])     # Output: apple\nprint(fruits[-1])    # Output: date\nprint(fruits[1:3])   # Output: ['banana', 'cherry']\n\n# Modifying lists\nfruits[1] = \"blueberry\"  # Change an element\nfruits.append(\"elderberry\")  # Add to end\nfruits.insert(1, \"apricot\")  # Insert at position\nfruits.remove(\"cherry\")  # Remove by value\npopped = fruits.pop(1)  # Remove by index and return")
    
    add_content_slide(prs, 
                     "List Operations", 
                     [
                         "• len(): Get the length of a list",
                         "• count(): Count occurrences of an element",
                         "• index(): Find the index of an element",
                         "• sort(): Sort the list in-place",
                         "• reverse(): Reverse the list in-place",
                         "• List comprehensions: Concise way to create lists"
                     ],
                     code="# List operations\nnumbers = [3, 1, 4, 1, 5, 9, 2]\nprint(len(numbers))       # Output: 7\nprint(numbers.count(1))   # Output: 2\nprint(numbers.index(5))   # Output: 4\n\nnumbers.sort()\nprint(numbers)  # Output: [1, 1, 2, 3, 4, 5, 9]\n\n# List comprehension\nsquares = [x**2 for x in range(1, 6)]\nprint(squares)  # Output: [1, 4, 9, 16, 25]\n\nevens = [x for x in range(1, 11) if x % 2 == 0]\nprint(evens)  # Output: [2, 4, 6, 8, 10]")
    
    # Tuples
    add_content_slide(prs, 
                     "Tuples", 
                     [
                         "• Ordered, immutable collections",
                         "• Created with parentheses ()",
                         "• Cannot be changed after creation",
                         "• Faster than lists",
                         "• Can be used as dictionary keys",
                         "• Good for data that shouldn't change"
                     ],
                     code="# Creating tuples\nempty_tuple = ()\nnumbers = (1, 2, 3, 4, 5)\nmixed = (1, \"hello\", 3.14, True)\nsingle = (1,)  # Note the comma for single-element tuples\n\n# Accessing elements (same as lists)\ncoordinates = (10, 20, 30, 40, 50)\nprint(coordinates[0])    # Output: 10\nprint(coordinates[-1])   # Output: 50\nprint(coordinates[1:3])  # Output: (20, 30)\n\n# Tuple operations\nprint(len(coordinates))  # Output: 5\nprint(coordinates.count(30))  # Output: 1\nprint(coordinates.index(40))  # Output: 3\n\n# Tuple unpacking\npoint = (3, 4)\nx, y = point\nprint(f\"X: {x}, Y: {y}\")  # Output: X: 3, Y: 4")
    
    # Sets
    add_content_slide(prs, 
                     "Sets", 
                     [
                         "• Unordered collections of unique elements",
                         "• Created with curly braces {} or set()",
                         "• No duplicate elements",
                         "• Fast membership testing",
                         "• Support mathematical set operations",
                         "• Cannot be indexed (no order)"
                     ],
                     code="# Creating sets\nempty_set = set()  # Not {} which creates an empty dictionary\nnumbers = {1, 2, 3, 4, 5}\nmixed = {1, \"hello\", 3.14, True}\n\n# Remove duplicates from a list\nnumbers_list = [1, 2, 2, 3, 4, 4, 5]\nunique_numbers = set(numbers_list)\nprint(unique_numbers)  # Output: {1, 2, 3, 4, 5}\n\n# Set operations\nset1 = {1, 2, 3, 4, 5}\nset2 = {4, 5, 6, 7, 8}\n\n# Union (all elements from both sets)\nprint(set1 | set2)  # Output: {1, 2, 3, 4, 5, 6, 7, 8}\n\n# Intersection (elements in both sets)\nprint(set1 & set2)  # Output: {4, 5}\n\n# Difference (elements in set1 but not in set2)\nprint(set1 - set2)  # Output: {1, 2, 3}")
    
    # Dictionaries
    add_content_slide(prs, 
                     "Dictionaries", 
                     [
                         "• Unordered collections of key-value pairs",
                         "• Created with curly braces {} and colons",
                         "• Keys must be unique and immutable",
                         "• Values can be any type",
                         "• Fast lookup by key",
                         "• Extremely useful for representing real-world data"
                     ],
                     code="# Creating dictionaries\nempty_dict = {}\nperson = {\n    \"name\": \"Alice\",\n    \"age\": 30,\n    \"city\": \"New York\"\n}\n\n# Accessing values\nprint(person[\"name\"])  # Output: Alice\nprint(person.get(\"age\"))  # Output: 30\nprint(person.get(\"country\", \"Unknown\"))  # Output: Unknown\n\n# Modifying dictionaries\nperson[\"email\"] = \"alice@example.com\"  # Add new key-value pair\nperson[\"age\"] = 31  # Update existing value\nremoved_value = person.pop(\"city\")  # Remove and return value\nprint(person)  # Output: {'name': 'Alice', 'age': 31, 'email': 'alice@example.com'}")
    
    add_content_slide(prs, 
                     "Dictionary Operations", 
                     [
                         "• keys(): Get all keys",
                         "• values(): Get all values",
                         "• items(): Get all key-value pairs as tuples",
                         "• update(): Update with another dictionary",
                         "• Nested dictionaries: Dictionaries inside dictionaries",
                         "• Dictionary comprehensions: Concise way to create dictionaries"
                     ],
                     code="# Dictionary operations\nperson = {\"name\": \"Alice\", \"age\": 30, \"city\": \"New York\"}\n\n# Get all keys, values, and items\nprint(list(person.keys()))    # Output: ['name', 'age', 'city']\nprint(list(person.values()))  # Output: ['Alice', 30, 'New York']\nprint(list(person.items()))   # Output: [('name', 'Alice'), ('age', 30), ('city', 'New York')]\n\n# Check if key exists\nprint(\"name\" in person)  # Output: True\nprint(\"email\" in person)  # Output: False\n\n# Dictionary comprehension\nsquares = {x: x**2 for x in range(1, 6)}\nprint(squares)  # Output: {1: 1, 2: 4, 3: 9, 4: 16, 5: 25}")
    
    # Real-life Project section
    add_section_slide(prs, "Real-life Project: Student Grading System")
    
    add_content_slide(prs, 
                     "Student Grading System", 
                     [
                         "• Comprehensive project applying functions, loops, and data structures",
                         "• Features:",
                         "  - Adding students",
                         "  - Recording grades",
                         "  - Calculating averages",
                         "  - Generating reports",
                         "  - Exporting/importing data",
                         "",
                         "• Demonstrates practical application of Python concepts"
                     ])
    
    # Project Structure slide - fixed the syntax error
    content = [
        "• Data structure: Nested dictionaries",
        "• Functions for each feature",
        "• Loops for processing data",
        "• Input validation",
        "• File I/O for data persistence",
        "",
        "• Key functions:",
        "  - add_student()",
        "  - record_grade()",
        "  - view_student_grades()",
        "  - class_summary()",
        "  - subject_summary()"
    ]
    
    code_example = "# Data structure\nstudents = {\n    \"student_id\": {\n        'name': \"Student Name\",\n        'grades': {\n            'Subject1': [grade1, grade2, grade3],\n            'Subject2': [grade1, grade2, grade3]\n        }\n    }\n}\n\n# Example function\ndef calculate_average(grades):\n    # Calculate the average of a list of grades\n    if not grades:\n        return 0\n    return sum(grades) / len(grades)"
    
    add_content_slide(prs, "Project Structure", content, code_example)
    
    add_content_slide(prs, 
                     "Project Demo", 
                     [
                         "• Let's explore the Student Grading System",
                         "• Main features:",
                         "  - Adding students with unique IDs",
                         "  - Recording grades for different subjects",
                         "  - Viewing individual student performance",
                         "  - Generating class summary reports",
                         "  - Analyzing subject-specific performance",
                         "",
                         "• The project demonstrates how to apply Python concepts to solve real-world problems"
                     ])
    
    # Conclusion slide
    add_content_slide(prs, 
                     "Conclusion", 
                     [
                         "• Functions help organize code and make it reusable",
                         "• Loops allow you to repeat code execution efficiently",
                         "• Data structures provide ways to organize and manage data:",
                         "  - Lists: Ordered, mutable collections",
                         "  - Tuples: Ordered, immutable collections",
                         "  - Sets: Unordered collections of unique elements",
                         "  - Dictionaries: Key-value pairs",
                         "",
                         "• These concepts are fundamental building blocks for Python programming",
                         "• The Student Grading System project shows how to apply these concepts in practice"
                     ])
    
    # Save the presentation
    prs.save('/home/ubuntu/python_day2_lecture/presentation/Python_Day2_Lecture.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
