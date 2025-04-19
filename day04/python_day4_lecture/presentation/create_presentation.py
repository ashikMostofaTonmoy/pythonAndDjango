"""
Create PowerPoint presentation for Python Day 4 lecture covering String, JSON, Libraries, Iterators, Modules, Dates, Math, Regex, Try/Except, and String formatting
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    """Add a title slide with the given title and subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_section_slide(prs, title):
    """Add a section slide with the given title."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_content_slide(prs, title, content, code=None):
    """Add a content slide with the given title, content, and optional code."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        
        # Check if this is a bullet point
        if line.startswith('•'):
            p.level = 1
    
    # Add code if provided
    if code:
        left = Inches(0.5)
        top = Inches(4)
        width = Inches(9)
        height = Inches(2.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.name = 'Courier New'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        # Add a light gray background to the code box
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide

def create_presentation():
    """Create the PowerPoint presentation for Python Day 4 lecture."""
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, 
                   "Python Day 4: String, JSON, Libraries, and Error Handling", 
                   "Advanced Python Concepts for Beginners")
    
    # Introduction slide
    add_content_slide(prs, 
                     "Today's Topics", 
                     [
                         "• Strings and String Formatting",
                         "• Iterators and Modules",
                         "• Date and Time",
                         "• Math Module",
                         "• JSON in Python",
                         "• Regular Expressions (Regex)",
                         "• Exception Handling (Try and Except)",
                         "• Libraries and Modules",
                         "",
                         "• Real-life Projects:",
                         "  - Task Scheduler with Dates and JSON Data",
                         "  - Regex-Based Log File Analyzer"
                     ])
    
    # Strings and String Formatting section
    add_section_slide(prs, "Strings and String Formatting")
    
    add_content_slide(prs, 
                     "String Basics", 
                     [
                         "• Strings are sequences of characters",
                         "• Created with single, double, or triple quotes",
                         "• Immutable (cannot be changed after creation)",
                         "• Accessed by index (zero-based)",
                         "• Sliced using square brackets [start:end:step]"
                     ],
                     code="# String creation\nsingle_quoted = 'Hello, World!'\ndouble_quoted = \"Hello, World!\"\ntriple_quoted = '''This is a multi-line\nstring that spans\nmultiple lines.'''\n\n# String length\nprint(len(single_quoted))  # Output: 13\n\n# Accessing characters\nprint(single_quoted[0])    # Output: H\nprint(single_quoted[-1])   # Output: !\n\n# Slicing\nprint(single_quoted[0:5])  # Output: Hello")
    
    add_content_slide(prs, 
                     "String Methods", 
                     [
                         "• Python provides many built-in methods for string manipulation",
                         "• Case conversion: upper(), lower(), capitalize(), title()",
                         "• Finding and counting: find(), count()",
                         "• Checking properties: isalpha(), isdigit(), isalnum(), isspace()",
                         "• Stripping whitespace: strip(), lstrip(), rstrip()",
                         "• Replacing text: replace()",
                         "• Splitting and joining: split(), join()"
                     ],
                     code="text = \"hello, world!\"\n\n# Case conversion\nprint(text.upper())        # Output: HELLO, WORLD!\nprint(text.capitalize())   # Output: Hello, world!\n\n# Finding and counting\nprint(text.find(\"world\"))  # Output: 7\nprint(text.count(\"l\"))     # Output: 3\n\n# Splitting and joining\nwords = text.split(\", \")   # Output: ['hello', 'world!']\nprint(\", \".join([\"hello\", \"Python\"]))  # Output: hello, Python")
    
    add_content_slide(prs, 
                     "String Formatting", 
                     [
                         "• Python offers several ways to format strings:",
                         "  1. % Operator (old style)",
                         "  2. format() Method",
                         "  3. f-strings (Formatted String Literals, Python 3.6+)",
                         "",
                         "• f-strings are the most modern and readable approach",
                         "• Allow expressions inside curly braces {}",
                         "• Support formatting specifiers for numbers, dates, etc."
                     ],
                     code="name = \"Alice\"\nage = 30\n\n# 1. % Operator (old style)\nprint(\"My name is %s and I am %d years old.\" % (name, age))\n\n# 2. format() Method\nprint(\"My name is {} and I am {} years old.\".format(name, age))\nprint(\"My name is {name} and I am {age} years old.\".format(name=name, age=age))\n\n# 3. f-strings (Python 3.6+)\nprint(f\"My name is {name} and I am {age} years old.\")\nprint(f\"In 5 years, I will be {age + 5} years old.\")\n\n# Formatting numbers\npi = 3.14159\nprint(f\"Pi is approximately {pi:.2f}\")")
    
    # Iterators and Modules section
    add_section_slide(prs, "Iterators and Modules")
    
    add_content_slide(prs, 
                     "Iterators", 
                     [
                         "• An iterator is an object that can be iterated upon",
                         "• Implements the iterator protocol (__iter__() and __next__())",
                         "• Used in for loops, list comprehensions, etc.",
                         "• Built-in iter() and next() functions",
                         "• StopIteration exception is raised when there are no more items"
                     ],
                     code="# Using an iterator\nmy_list = [1, 2, 3, 4, 5]\nmy_iter = iter(my_list)  # Get an iterator\n\nprint(next(my_iter))  # Output: 1\nprint(next(my_iter))  # Output: 2\nprint(next(my_iter))  # Output: 3\n\n# Creating a custom iterator\nclass CountDown:\n    def __init__(self, start):\n        self.start = start\n        \n    def __iter__(self):\n        return self\n        \n    def __next__(self):\n        if self.start <= 0:\n            raise StopIteration\n        self.start -= 1\n        return self.start + 1")
    
    add_content_slide(prs, 
                     "Generators", 
                     [
                         "• Generators are a simple way of creating iterators",
                         "• Created using functions with the yield statement",
                         "• Generate values on the fly (lazy evaluation)",
                         "• More memory-efficient than storing all values in a list",
                         "• Generator expressions are similar to list comprehensions"
                     ],
                     code="# Generator function\ndef countdown(n):\n    while n > 0:\n        yield n\n        n -= 1\n\n# Using the generator\nfor num in countdown(5):\n    print(num)  # Output: 5, 4, 3, 2, 1\n\n# Generator expression\nsquares = (x**2 for x in range(1, 6))\nfor square in squares:\n    print(square)  # Output: 1, 4, 9, 16, 25")
    
    add_content_slide(prs, 
                     "Modules", 
                     [
                         "• A module is a file containing Python definitions and statements",
                         "• Used to organize code into reusable components",
                         "• Import using the import statement",
                         "• Can import the entire module or specific functions/classes",
                         "• Can use aliases to avoid naming conflicts",
                         "• dir() function shows available attributes and methods"
                     ],
                     code="# Import the entire module\nimport math\nprint(math.sqrt(16))  # Output: 4.0\n\n# Import specific functions\nfrom math import sqrt, pi\nprint(sqrt(16))       # Output: 4.0\n\n# Import with an alias\nimport math as m\nprint(m.sqrt(16))     # Output: 4.0\n\n# Creating your own module (in mymodule.py)\ndef greet(name):\n    return f\"Hello, {name}!\"\n\n# Using your module\nimport mymodule\nprint(mymodule.greet(\"Alice\"))  # Output: Hello, Alice!")
    
    # Date and Time section
    add_section_slide(prs, "Date and Time")
    
    add_content_slide(prs, 
                     "Basic Date and Time Operations", 
                     [
                         "• Python's datetime module provides classes for manipulating dates and times",
                         "• Key classes: datetime, date, time, timedelta",
                         "• Creating date and time objects",
                         "• Accessing components (year, month, day, hour, etc.)",
                         "• Getting the current date and time"
                     ],
                     code="import datetime\n\n# Current date and time\nnow = datetime.datetime.now()\nprint(now)  # Output: 2023-04-19 13:30:45.123456\n\n# Creating date objects\ndate1 = datetime.date(2023, 4, 19)\nprint(date1)  # Output: 2023-04-19\n\n# Creating time objects\ntime1 = datetime.time(13, 30, 45)\nprint(time1)  # Output: 13:30:45\n\n# Date components\ntoday = datetime.date.today()\nprint(today.year)    # Output: 2023\nprint(today.month)   # Output: 4\nprint(today.day)     # Output: 19")
    
    add_content_slide(prs, 
                     "Date Formatting and Parsing", 
                     [
                         "• Converting dates to strings using strftime()",
                         "• Parsing strings into dates using strptime()",
                         "• Format codes: %Y (year), %m (month), %d (day), etc.",
                         "• Time format codes: %H (hour), %M (minute), %S (second), etc.",
                         "• Locale-specific formatting (day and month names)"
                     ],
                     code="import datetime\n\nnow = datetime.datetime.now()\n\n# Format date using strftime()\nprint(now.strftime(\"%Y-%m-%d\"))  # Output: 2023-04-19\nprint(now.strftime(\"%d/%m/%Y\"))  # Output: 19/04/2023\nprint(now.strftime(\"%B %d, %Y\")) # Output: April 19, 2023\nprint(now.strftime(\"%H:%M:%S\"))  # Output: 13:30:45\n\n# Parse date string using strptime()\ndate_string = \"19 April, 2023\"\ndate_object = datetime.datetime.strptime(date_string, \"%d %B, %Y\")\nprint(date_object)  # Output: 2023-04-19 00:00:00")
    
    add_content_slide(prs, 
                     "Date Arithmetic", 
                     [
                         "• Adding or subtracting days, hours, minutes, etc.",
                         "• Using timedelta objects for date arithmetic",
                         "• Calculating the difference between dates",
                         "• Comparing dates (>, <, ==, etc.)",
                         "• Working with time zones (pytz library)"
                     ],
                     code="import datetime\n\ntoday = datetime.date.today()\n\n# Adding days\ntomorrow = today + datetime.timedelta(days=1)\nprint(tomorrow)  # Output: 2023-04-20\n\n# Subtracting days\nyesterday = today - datetime.timedelta(days=1)\nprint(yesterday)  # Output: 2023-04-18\n\n# Difference between dates\ndate1 = datetime.date(2023, 4, 19)\ndate2 = datetime.date(2023, 5, 1)\ndelta = date2 - date1\nprint(delta.days)  # Output: 12")
    
    # Math Module section
    add_section_slide(prs, "Math Module")
    
    add_content_slide(prs, 
                     "Basic Math Functions", 
                     [
                         "• Python's math module provides access to mathematical functions",
                         "• Constants: pi, e, inf, nan",
                         "• Rounding functions: ceil(), floor(), trunc()",
                         "• Power and logarithmic functions: pow(), sqrt(), log()",
                         "• Trigonometric functions: sin(), cos(), tan()",
                         "• Conversion functions: degrees(), radians()"
                     ],
                     code="import math\n\n# Constants\nprint(math.pi)       # Output: 3.141592653589793\nprint(math.e)        # Output: 2.718281828459045\n\n# Rounding functions\nprint(math.ceil(4.2))   # Output: 5 (rounds up)\nprint(math.floor(4.8))  # Output: 4 (rounds down)\n\n# Power and logarithmic functions\nprint(math.pow(2, 3))   # Output: 8.0 (2^3)\nprint(math.sqrt(16))    # Output: 4.0 (square root)\nprint(math.log10(100))  # Output: 2.0 (log base 10)")
    
    add_content_slide(prs, 
                     "Statistical Functions and Random Module", 
                     [
                         "• statistics module for statistical functions",
                         "• mean(), median(), mode(), stdev(), variance()",
                         "• random module for random number generation",
                         "• random(), uniform(), randint(), choice(), sample()",
                         "• shuffle() for randomizing sequences"
                     ],
                     code="import statistics\nimport random\n\ndata = [1, 2, 3, 4, 5, 5, 6, 7, 8, 9]\n\n# Basic statistics\nprint(statistics.mean(data))     # Output: 5.0 (average)\nprint(statistics.median(data))   # Output: 5.0 (middle value)\nprint(statistics.mode(data))     # Output: 5 (most common value)\n\n# Random numbers\nprint(random.random())  # Output: 0.123456789... (0 to 1)\nprint(random.randint(1, 10))  # Output: 7 (1 to 10)\nprint(random.choice(['apple', 'banana', 'cherry']))  # Output: 'banana'")
    
    # JSON section
    add_section_slide(prs, "JSON in Python")
    
    add_content_slide(prs, 
                     "JSON Basics", 
                     [
                         "• JSON (JavaScript Object Notation) is a lightweight data interchange format",
                         "• Human-readable and machine-parsable",
                         "• Common format for web APIs and configuration files",
                         "• Python's json module for encoding and decoding JSON",
                         "• JSON data types map to Python types (dict, list, str, int, float, bool, None)"
                     ],
                     code="import json\n\n# Python dictionary\nperson = {\n    \"name\": \"John Doe\",\n    \"age\": 30,\n    \"city\": \"New York\",\n    \"languages\": [\"Python\", \"JavaScript\", \"C++\"],\n    \"is_employee\": True\n}\n\n# JSON data types and their Python equivalents\n# object -> dict\n# array -> list\n# string -> str\n# number -> int/float\n# true/false -> True/False\n# null -> None")
    
    add_content_slide(prs, 
                     "Converting Python Objects to JSON", 
                     [
                         "• json.dumps() converts Python objects to JSON strings",
                         "• json.dump() writes Python objects to a JSON file",
                         "• Formatting options: indent, sort_keys",
                         "• Handling custom objects with default parameter or JSONEncoder"
                     ],
                     code="import json\n\nperson = {\n    \"name\": \"John Doe\",\n    \"age\": 30,\n    \"city\": \"New York\"\n}\n\n# Convert Python object to JSON string\njson_string = json.dumps(person)\nprint(json_string)\n# Output: {\"name\": \"John Doe\", \"age\": 30, \"city\": \"New York\"}\n\n# Pretty print with indentation\njson_formatted = json.dumps(person, indent=4)\nprint(json_formatted)\n\n# Writing JSON to a file\nwith open(\"person.json\", \"w\") as file:\n    json.dump(person, file, indent=4)")
    
    add_content_slide(prs, 
                     "Converting JSON to Python Objects", 
                     [
                         "• json.loads() converts JSON strings to Python objects",
                         "• json.load() reads JSON from a file into Python objects",
                         "• Handling custom types with object_hook parameter",
                         "• Error handling for invalid JSON"
                     ],
                     code="import json\n\n# JSON string\njson_string = '{\"name\": \"Jane Smith\", \"age\": 25, \"city\": \"London\"}'\n\n# Convert JSON string to Python object\nperson = json.loads(json_string)\nprint(person)  # Output: {'name': 'Jane Smith', 'age': 25, 'city': 'London'}\n\n# Access dictionary values\nprint(person[\"name\"])  # Output: Jane Smith\n\n# Reading JSON from a file\nwith open(\"person.json\", \"r\") as file:\n    loaded_person = json.load(file)\n    print(loaded_person)")
    
    # Regular Expressions section
    add_section_slide(prs, "Regular Expressions (Regex)")
    
    add_content_slide(prs, 
                     "Regex Basics", 
                     [
                         "• Regular expressions are patterns used to match character combinations in strings",
                         "• Python's re module for working with regular expressions",
                         "• Basic functions: search(), match(), findall(), sub()",
                         "• Compile patterns for reuse with re.compile()",
                         "• Raw strings (r'pattern') to avoid escaping backslashes"
                     ],
                     code="import re\n\ntext = \"The quick brown fox jumps over the lazy dog.\"\n\n# Search for a pattern\nmatch = re.search(r\"fox\", text)\nif match:\n    print(\"Pattern found:\", match.group())  # Output: Pattern found: fox\n    print(\"Position:\", match.start())       # Output: Position: 16\n\n# Find all occurrences\nmatches = re.findall(r\"the\", text, re.IGNORECASE)\nprint(matches)  # Output: ['The', 'the']\n\n# Replace pattern\nnew_text = re.sub(r\"fox\", \"cat\", text)\nprint(new_text)  # Output: The quick brown cat jumps over the lazy dog.")
    
    add_content_slide(prs, 
                     "Metacharacters and Special Sequences", 
                     [
                         "• . (dot) - Matches any character except newline",
                         "• ^ - Matches start of string",
                         "• $ - Matches end of string",
                         "• * - Matches 0 or more occurrences",
                         "• + - Matches 1 or more occurrences",
                         "• ? - Matches 0 or 1 occurrence",
                         "• {} - Matches specified number of occurrences",
                         "• [] - Matches any character in the brackets",
                         "• | - Alternation (OR)",
                         "• () - Groups patterns"
                     ],
                     code="import re\n\n# . (dot) - Matches any character except newline\nprint(re.findall(r\"b.t\", \"bit bat but bet\"))  # Output: ['bit', 'bat', 'but', 'bet']\n\n# * (asterisk) - Matches 0 or more occurrences\nprint(re.findall(r\"ab*c\", \"ac abc abbc\"))  # Output: ['ac', 'abc', 'abbc']\n\n# [] (square brackets) - Matches any character in the brackets\nprint(re.findall(r\"[aeiou]\", \"apple\"))  # Output: ['a', 'e']\n\n# | (pipe) - Alternation (OR)\nprint(re.findall(r\"cat|dog\", \"I have a cat and a dog\"))  # Output: ['cat', 'dog']")
    
    add_content_slide(prs, 
                     "Special Sequences and Practical Examples", 
                     [
                         "• \\d - Matches any decimal digit (0-9)",
                         "• \\D - Matches any non-digit character",
                         "• \\w - Matches any alphanumeric character (a-z, A-Z, 0-9, _)",
                         "• \\W - Matches any non-alphanumeric character",
                         "• \\s - Matches any whitespace character",
                         "• \\S - Matches any non-whitespace character",
                         "• \\b - Matches word boundary"
                     ],
                     code="import re\n\n# Validating email addresses\ndef is_valid_email(email):\n    pattern = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\\.[a-zA-Z]{2,}$'\n    return bool(re.match(pattern, email))\n\nprint(is_valid_email(\"user@example.com\"))  # Output: True\nprint(is_valid_email(\"invalid-email\"))     # Output: False\n\n# Extracting phone numbers\ntext = \"Contact us at (123) 456-7890 or 987-654-3210\"\nphone_numbers = re.findall(r'\\(?\\d{3}\\)?[-.]?\\d{3}[-.]?\\d{4}', text)\nprint(phone_numbers)  # Output: ['(123) 456-7890', '987-654-3210']")
    
    # Exception Handling section
    add_section_slide(prs, "Exception Handling (Try and Except)")
    
    add_content_slide(prs, 
                     "Basic Exception Handling", 
                     [
                         "• Exceptions are errors that occur during program execution",
                         "• try-except blocks catch and handle exceptions",
                         "• Prevents program crashes and allows graceful error handling",
                         "• Can catch specific exception types or all exceptions",
                         "• Multiple except blocks for different exception types"
                     ],
                     code="# Without exception handling\n# x = 10 / 0  # This would raise a ZeroDivisionError and crash the program\n\n# With exception handling\ntry:\n    x = 10 / 0\nexcept ZeroDivisionError:\n    print(\"Error: Division by zero!\")  # Output: Error: Division by zero!\n\n# Handling multiple exceptions\ntry:\n    # This could raise different exceptions\n    num = int(input(\"Enter a number: \"))  # Let's say user enters \"abc\"\n    result = 10 / num\nexcept ValueError:\n    print(\"Error: Please enter a valid number!\")\nexcept ZeroDivisionError:\n    print(\"Error: Division by zero!\")")
    
    add_content_slide(prs, 
                     "The else and finally Clauses", 
                     [
                         "• else clause executes if no exceptions were raised",
                         "• finally clause always executes, regardless of exceptions",
                         "• Used for cleanup operations (closing files, releasing resources)",
                         "• Can be combined with try-except blocks",
                         "• Execution flow with try-except-else-finally"
                     ],
                     code="try:\n    num = int(input(\"Enter a number: \"))  # Let's say user enters \"5\"\n    result = 10 / num\nexcept ValueError:\n    print(\"Error: Please enter a valid number!\")\nexcept ZeroDivisionError:\n    print(\"Error: Division by zero!\")\nelse:\n    # This block executes if no exceptions were raised\n    print(f\"Result: {result}\")  # Output: Result: 2.0\nfinally:\n    # This block always executes\n    print(\"Execution completed.\")  # Output: Execution completed.")
    
    add_content_slide(prs, 
                     "Raising Exceptions and Custom Exceptions", 
                     [
                         "• raise statement to manually trigger exceptions",
                         "• Creating custom exception classes by inheriting from Exception",
                         "• Adding custom attributes and methods to exceptions",
                         "• Using custom exceptions for specific error cases",
                         "• Best practices for exception handling"
                     ],
                     code="# Raising exceptions\ndef validate_age(age):\n    if age < 0:\n        raise ValueError(\"Age cannot be negative\")\n    if age > 120:\n        raise ValueError(\"Age is too high\")\n    return age\n\n# Creating custom exceptions\nclass CustomError(Exception):\n    \"\"\"Base class for custom exceptions\"\"\"\n    pass\n\nclass ValueTooSmallError(CustomError):\n    \"\"\"Raised when the input value is too small\"\"\"\n    pass\n\nclass ValueTooLargeError(CustomError):\n    \"\"\"Raised when the input value is too large\"\"\"\n    pass")
    
    # Libraries and Modules section
    add_section_slide(prs, "Libraries and Modules")
    
    add_content_slide(prs, 
                     "Standard Library Modules", 
                     [
                         "• Python comes with a comprehensive standard library",
                         "• os - Operating system interface",
                         "• sys - System-specific parameters and functions",
                         "• collections - Specialized container datatypes",
                         "• itertools - Functions for efficient looping",
                         "• functools - Higher-order functions and operations on callable objects",
                         "• Many more modules for various tasks"
                     ],
                     code="# os - Operating system interface\nimport os\nprint(os.getcwd())  # Output: Current working directory\nprint(os.listdir())  # Output: List of files in current directory\n\n# collections - Specialized container datatypes\nfrom collections import Counter, defaultdict, namedtuple\n\n# Count occurrences of elements\ncounter = Counter(['apple', 'banana', 'apple', 'orange', 'banana', 'apple'])\nprint(counter)  # Output: Counter({'apple': 3, 'banana': 2, 'orange': 1})\n\n# Dictionary with default values\nfruit_count = defaultdict(int)  # Default value is 0\nfruit_count['apple'] += 1\nprint(fruit_count['apple'])    # Output: 1\nprint(fruit_count['banana'])   # Output: 0 (default value)")
    
    add_content_slide(prs, 
                     "Third-Party Libraries", 
                     [
                         "• Python's ecosystem includes thousands of third-party libraries",
                         "• requests - HTTP library for making requests",
                         "• pandas - Data analysis and manipulation",
                         "• numpy - Numerical computing",
                         "• matplotlib - Data visualization",
                         "• Installing packages with pip",
                         "• Virtual environments for isolated dependencies"
                     ],
                     code="# requests - HTTP library for making requests\n# pip install requests\nimport requests\nresponse = requests.get('https://api.github.com')\nprint(response.status_code)  # Output: 200\n\n# pandas - Data analysis and manipulation\n# pip install pandas\nimport pandas as pd\ndata = {'Name': ['John', 'Anna', 'Peter', 'Linda'],\n        'Age': [28, 24, 35, 32]}\ndf = pd.DataFrame(data)\nprint(df)\n\n# Virtual environments\n# python -m venv myenv\n# myenv\\Scripts\\activate (Windows)\n# source myenv/bin/activate (macOS/Linux)")
    
    # Real-life Projects section
    add_section_slide(prs, "Real-life Projects")
    
    add_content_slide(prs, 
                     "Task Scheduler with Dates and JSON Data", 
                     [
                         "• Command-line application for managing tasks with due dates",
                         "• Features:",
                         "  - Add tasks with titles, descriptions, and due dates",
                         "  - Save tasks to a JSON file",
                         "  - Load tasks from a JSON file",
                         "  - List tasks sorted by due date",
                         "  - Mark tasks as completed or pending",
                         "  - Filter tasks by status, due date, or date range",
                         "",
                         "• Demonstrates date handling, JSON processing, and exception handling"
                     ])
    
    add_content_slide(prs, 
                     "Task Scheduler Implementation", 
                     [
                         "• Object-oriented design with three main classes:",
                         "  - Task: Represents a single task with properties",
                         "  - TaskManager: Manages a collection of tasks",
                         "  - TaskSchedulerCLI: Provides a command-line interface",
                         "",
                         "• Key concepts demonstrated:",
                         "  - Date and time handling with datetime module",
                         "  - JSON serialization/deserialization",
                         "  - Exception handling for file operations",
                         "  - Command-line interface design"
                     ],
                     code="class Task:\n    def __init__(self, title, description, due_date, completed=False, task_id=None):\n        self.title = title\n        self.description = description\n        self.due_date = due_date\n        self.completed = completed\n        self.task_id = task_id\n        self.created_at = datetime.datetime.now()\n    \n    def to_dict(self):\n        return {\n            'task_id': self.task_id,\n            'title': self.title,\n            'description': self.description,\n            'due_date': self.due_date.isoformat(),\n            'completed': self.completed,\n            'created_at': self.created_at.isoformat()\n        }")
    
    add_content_slide(prs, 
                     "Regex-Based Log File Analyzer", 
                     [
                         "• Command-line application for analyzing log files",
                         "• Features:",
                         "  - Parse log files using regular expressions",
                         "  - Extract important information (timestamps, IP addresses, etc.)",
                         "  - Generate statistics and reports",
                         "  - Filter logs based on various criteria",
                         "  - Detect patterns or anomalies in the logs",
                         "",
                         "• Demonstrates regex usage, string manipulation, and data analysis"
                     ])
    
    add_content_slide(prs, 
                     "Log File Analyzer Implementation", 
                     [
                         "• Object-oriented design with four main classes:",
                         "  - LogEntry: Represents a single log entry",
                         "  - LogParser: Parses log files using regular expressions",
                         "  - LogAnalyzer: Analyzes log entries and generates reports",
                         "  - LogAnalyzerCLI: Provides a command-line interface",
                         "",
                         "• Key concepts demonstrated:",
                         "  - Regular expressions for pattern matching",
                         "  - String parsing and manipulation",
                         "  - Data analysis and reporting",
                         "  - Command-line argument parsing"
                     ],
                     code="class LogParser:\n    # Common log format pattern (Apache/Nginx)\n    COMMON_LOG_PATTERN = re.compile(\n        r'(\\d+\\.\\d+\\.\\d+\\.\\d+) - - \\[(.*?)\\] \"([A-Z]+) (.*?) HTTP/\\d\\.\\d\" (\\d+) (\\d+)'\n    )\n    \n    @classmethod\n    def parse_log_line(cls, log_line):\n        match = cls.COMMON_LOG_PATTERN.match(log_line)\n        if match:\n            ip_address, timestamp_str, method, path, status_code, response_size = match.groups()\n            # Process and return LogEntry object")
    
    # Conclusion slide
    add_content_slide(prs, 
                     "Conclusion", 
                     [
                         "• We've covered advanced Python concepts for beginners:",
                         "  - Strings and string formatting",
                         "  - Iterators and modules",
                         "  - Date and time handling",
                         "  - Math functions and random numbers",
                         "  - JSON data processing",
                         "  - Regular expressions",
                         "  - Exception handling",
                         "  - Libraries and modules",
                         "",
                         "• Applied these concepts in two real-life projects",
                         "• These skills will help you build more complex and robust Python applications"
                     ])
    
    # Save the presentation
    prs.save('/home/ubuntu/python_day4_lecture/presentation/Python_Day4_Lecture.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
