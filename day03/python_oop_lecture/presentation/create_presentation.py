"""
Create PowerPoint presentation for Python OOP lecture
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    """Add a title slide with the given title and subtitle."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_section_slide(prs, title):
    """Add a section slide with the given title."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    return slide

def add_content_slide(prs, title, content, code=None):
    """Add a content slide with the given title, content, and optional code."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    
    # Format title
    title_shape = slide.shapes.title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 112, 192)
    
    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for line in content:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(24)
        
        # Check if this is a bullet point
        if line.startswith('•'):
            p.level = 1
    
    # Add code if provided
    if code:
        left = Inches(0.5)
        top = Inches(4)
        width = Inches(9)
        height = Inches(2.5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        
        p = text_frame.add_paragraph()
        p.text = code
        p.font.name = 'Courier New'
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0, 0, 128)
        
        # Add a light gray background to the code box
        fill = textbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 240, 240)
    
    return slide

def create_presentation():
    """Create the PowerPoint presentation for Python OOP lecture."""
    prs = Presentation()
    
    # Set slide dimensions to 16:9 aspect ratio
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    add_title_slide(prs, 
                   "Object-Oriented Programming in Python", 
                   "Lambda Functions, Arrays, Classes, Objects, Inheritance, and Polymorphism")
    
    # Introduction slide
    add_content_slide(prs, 
                     "Introduction to OOP", 
                     [
                         "• Object-Oriented Programming (OOP) is a programming paradigm",
                         "• Uses objects and classes to structure code",
                         "• Makes complex code more manageable, reusable, and organized",
                         "• Python supports multiple programming paradigms, including OOP",
                         "",
                         "Today we'll cover:",
                         "• Lambda Functions",
                         "• Arrays (Lists in Python)",
                         "• Classes and Objects",
                         "• Inheritance",
                         "• Polymorphism"
                     ])
    
    # Lambda Functions section
    add_section_slide(prs, "Lambda Functions")
    
    add_content_slide(prs, 
                     "What are Lambda Functions?", 
                     [
                         "• Small, anonymous functions defined with the lambda keyword",
                         "• Can take any number of arguments but only have one expression",
                         "• Useful for short, simple functions that are used only once",
                         "",
                         "Syntax:",
                         "lambda arguments: expression"
                     ])
    
    add_content_slide(prs, 
                     "Lambda Function Examples", 
                     [
                         "Example 1: Simple lambda function",
                         "",
                         "Example 2: Lambda with sorted()",
                         "",
                         "Example 3: Lambda with filter()"
                     ],
                     code="""# Example 1: Simple lambda function
add = lambda x, y: x + y
print(add(5, 3))  # Output: 8

# Example 2: Lambda with sorted()
pairs = [(1, 'one'), (3, 'three'), (2, 'two')]
sorted_pairs = sorted(pairs, key=lambda pair: pair[1])
print(sorted_pairs)  # [(4, 'four'), (1, 'one'), (3, 'three'), (2, 'two')]

# Example 3: Lambda with filter()
numbers = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10]
even_numbers = list(filter(lambda x: x % 2 == 0, numbers))
print(even_numbers)  # [2, 4, 6, 8, 10]""")
    
    # Arrays (Lists) section
    add_section_slide(prs, "Arrays (Lists in Python)")
    
    add_content_slide(prs, 
                     "Python Lists", 
                     [
                         "• In Python, we typically use lists as arrays",
                         "• Lists are ordered, mutable collections",
                         "• Can contain elements of different types",
                         "• Dynamic sizing (no need to declare size in advance)",
                         "",
                         "Creating Lists:",
                         "• Empty list: []",
                         "• List with elements: [1, 2, 3]",
                         "• Mixed types: [1, 'hello', 3.14, True]"
                     ])
    
    add_content_slide(prs, 
                     "List Operations", 
                     [
                         "• Accessing elements with indexing",
                         "• Slicing to get portions of lists",
                         "• Adding elements with append(), insert(), extend()",
                         "• Removing elements with remove(), pop()",
                         "• Other operations: len(), count(), sort(), reverse()"
                     ],
                     code="""# Accessing elements
fruits = ["apple", "banana", "cherry", "date"]
print(fruits[0])    # Output: apple
print(fruits[-1])   # Output: date

# Slicing
print(fruits[1:3])  # Output: ['banana', 'cherry']

# Adding elements
fruits.append("elderberry")
fruits.insert(1, "apricot")
print(fruits)  # ['apple', 'apricot', 'banana', 'cherry', 'date', 'elderberry']

# Removing elements
fruits.remove("banana")
popped = fruits.pop(1)  # Removes and returns 'apricot'""")
    
    add_content_slide(prs, 
                     "List Comprehensions", 
                     [
                         "• Concise way to create lists",
                         "• More readable and often faster than loops",
                         "• Can include conditions and nested loops",
                         "",
                         "Syntax:",
                         "[expression for item in iterable if condition]"
                     ],
                     code="""# Create a list of squares
squares = [x**2 for x in range(1, 6)]
print(squares)  # Output: [1, 4, 9, 16, 25]

# Create a list of even numbers
evens = [x for x in range(1, 11) if x % 2 == 0]
print(evens)  # Output: [2, 4, 6, 8, 10]

# Create a matrix (list of lists)
matrix = [[i * j for j in range(1, 4)] for i in range(1, 4)]
print(matrix)  # Output: [[1, 2, 3], [2, 4, 6], [3, 6, 9]]""")
    
    # Classes and Objects section
    add_section_slide(prs, "Classes and Objects")
    
    add_content_slide(prs, 
                     "Classes and Objects", 
                     [
                         "• Classes are blueprints for creating objects",
                         "• Define attributes (data) and methods (functions)",
                         "• Objects are instances of classes",
                         "",
                         "Class components:",
                         "• __init__: Constructor method",
                         "• self: Reference to the instance",
                         "• Attributes: Variables that belong to the class",
                         "• Methods: Functions that belong to the class"
                     ])
    
    add_content_slide(prs, 
                     "Defining a Class", 
                     [
                         "• Use the class keyword to define a class",
                         "• Constructor method __init__ initializes new objects",
                         "• self parameter refers to the instance being created",
                         "• Instance methods operate on the object's data"
                     ],
                     code="""class Dog:
    # Class attribute (shared by all instances)
    species = "Canis familiaris"
    
    # Initializer / Constructor
    def __init__(self, name, age):
        # Instance attributes (unique to each instance)
        self.name = name
        self.age = age
    
    # Instance method
    def description(self):
        return f"{self.name} is {self.age} years old"
    
    # Another instance method
    def speak(self, sound):
        return f"{self.name} says {sound}" """)
    
    add_content_slide(prs, 
                     "Creating and Using Objects", 
                     [
                         "• Create objects by calling the class name as a function",
                         "• Access attributes with dot notation: object.attribute",
                         "• Call methods with dot notation: object.method()",
                         "• Each object has its own instance attributes",
                         "• All objects share class attributes"
                     ],
                     code="""# Create Dog objects
buddy = Dog("Buddy", 9)
miles = Dog("Miles", 4)

# Access attributes
print(buddy.name)  # Output: Buddy
print(miles.age)   # Output: 4

# Access class attribute
print(buddy.species)  # Output: Canis familiaris

# Call methods
print(buddy.description())  # Output: Buddy is 9 years old
print(miles.speak("Woof"))  # Output: Miles says Woof""")
    
    # Inheritance section
    add_section_slide(prs, "Inheritance")
    
    add_content_slide(prs, 
                     "Inheritance Basics", 
                     [
                         "• Inheritance allows a class to inherit attributes and methods from another class",
                         "• Parent/Base class: The original class",
                         "• Child/Derived class: The class that inherits",
                         "• Promotes code reuse and establishes relationships between classes",
                         "",
                         "Syntax:",
                         "class ChildClass(ParentClass):"
                     ])
    
    add_content_slide(prs, 
                     "Inheritance Example", 
                     [
                         "• Child class inherits all attributes and methods from parent",
                         "• Child class can override parent methods",
                         "• Child class can add new attributes and methods",
                         "• super() function calls methods from the parent class"
                     ],
                     code="""# Parent class
class Animal:
    def __init__(self, name, species):
        self.name = name
        self.species = species
    
    def make_sound(self):
        return "Some generic animal sound"

# Child class
class Dog(Animal):
    def __init__(self, name, breed):
        # Call parent class's __init__ method
        super().__init__(name, "Canis familiaris")
        self.breed = breed
    
    # Override parent method
    def make_sound(self):
        return "Woof!"
    
    # Add new method
    def wag_tail(self):
        return f"{self.name} wags tail" """)
    
    add_content_slide(prs, 
                     "Using Inheritance", 
                     [
                         "• Create objects of the child class",
                         "• Access inherited attributes and methods",
                         "• Access overridden methods",
                         "• Access child-specific methods"
                     ],
                     code="""# Create objects
generic_animal = Animal("Generic", "Animal species")
buddy = Dog("Buddy", "Golden Retriever")

# Use parent methods
print(generic_animal.make_sound())  # Some generic animal sound
print(buddy.make_sound())           # Woof!

# Access inherited attributes
print(buddy.name)    # Buddy
print(buddy.species) # Canis familiaris

# Use child-specific method
print(buddy.wag_tail())  # Buddy wags tail""")
    
    # Polymorphism section
    add_section_slide(prs, "Polymorphism")
    
    add_content_slide(prs, 
                     "Polymorphism Basics", 
                     [
                         "• Polymorphism: 'many forms'",
                         "• Allows objects of different classes to be treated as objects of a common base class",
                         "• Enables using a single interface with different underlying forms",
                         "• Two main types in Python:",
                         "  - Method overriding (inheritance-based)",
                         "  - Duck typing ('if it walks like a duck and quacks like a duck...')"
                     ])
    
    add_content_slide(prs, 
                     "Polymorphism Example", 
                     [
                         "• Different classes implement the same method name",
                         "• Each class provides its own implementation",
                         "• Code can work with different classes through a common interface",
                         "• Makes code more flexible and extensible"
                     ],
                     code="""class Animal:
    def speak(self):
        pass

class Dog(Animal):
    def speak(self):
        return "Woof!"

class Cat(Animal):
    def speak(self):
        return "Meow!"

class Duck(Animal):
    def speak(self):
        return "Quack!"

# Function that works with any Animal
def animal_sound(animal):
    return animal.speak()

# Polymorphic behavior
print(animal_sound(Dog()))   # Output: Woof!
print(animal_sound(Cat()))   # Output: Meow!
print(animal_sound(Duck()))  # Output: Quack!""")
    
    add_content_slide(prs, 
                     "Duck Typing", 
                     [
                         "• Python's approach to polymorphism",
                         "• Objects don't need to be from the same inheritance hierarchy",
                         "• Only need to implement the required methods",
                         "• 'If it walks like a duck and quacks like a duck, it's a duck'",
                         "• Focuses on behavior rather than type"
                     ],
                     code="""class Dog:
    def speak(self):
        return "Woof!"

class Cat:
    def speak(self):
        return "Meow!"

class Person:
    def speak(self):
        return "Hello!"

# Function that works with anything that has a speak method
def make_speak(entity):
    return entity.speak()

# Polymorphic behavior without inheritance
print(make_speak(Dog()))     # Output: Woof!
print(make_speak(Cat()))     # Output: Meow!
print(make_speak(Person()))  # Output: Hello!""")
    
    # Real-life Project section
    add_section_slide(prs, "Real-life Project: Online Bookstore")
    
    add_content_slide(prs, 
                     "Online Bookstore Management System", 
                     [
                         "• Comprehensive project applying OOP concepts",
                         "• Features:",
                         "  - User authentication",
                         "  - Inventory management",
                         "  - Shopping cart functionality",
                         "  - Order processing",
                         "",
                         "• Project structure:",
                         "  - Models: Book, User, Inventory",
                         "  - Controllers: Auth, Book, Cart",
                         "  - Views: Command Line Interface"
                     ])
    
    add_content_slide(prs, 
                     "Project Demo", 
                     [
                         "• Let's explore the project structure",
                         "• Key classes and their relationships:",
                         "  - Book: Represents books in inventory",
                         "  - User: Handles authentication and cart",
                         "  - Inventory: Manages collection of books",
                         "  - Controllers: Implement business logic",
                         "",
                         "• OOP concepts demonstrated:",
                         "  - Classes and Objects",
                         "  - Encapsulation",
                         "  - Inheritance (structure allows for it)",
                         "  - Polymorphism (consistent interfaces)"
                     ])
    
    # Conclusion slide
    add_content_slide(prs, 
                     "Conclusion", 
                     [
                         "• OOP in Python provides powerful tools for organizing code",
                         "• Lambda functions offer concise ways to create small functions",
                         "• Lists (arrays) are versatile data structures",
                         "• Classes define blueprints for objects, encapsulating data and behavior",
                         "• Inheritance allows for code reuse and establishing relationships",
                         "• Polymorphism enables flexibility and extensibility",
                         "",
                         "• These concepts help create more maintainable, reusable, and scalable applications"
                     ])
    
    # Save the presentation
    prs.save('/home/ubuntu/python_oop_lecture/presentation/Python_OOP_Lecture.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
